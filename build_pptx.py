"""Generate presentation.pptx (29 slides, native editable) from images/ + the
structure of presentation.md. No browser, no Marp, no internet required.

Run:
    pip install -r requirements.txt
    python build_pptx.py
"""

from __future__ import annotations

from pathlib import Path
from typing import Iterable, Sequence

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

COLOR_PRIMARY = RGBColor(0x1F, 0x3A, 0x5F)
COLOR_PRIMARY_LIGHT = RGBColor(0xE6, 0xEC, 0xF4)
COLOR_DIVIDER_BG = RGBColor(0x1F, 0x3A, 0x5F)
COLOR_PLACEHOLDER = RGBColor(0xB0, 0x00, 0x00)
COLOR_TITLE_BAR_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BODY_TEXT = RGBColor(0x22, 0x22, 0x22)
COLOR_SUBTLE = RGBColor(0x55, 0x55, 0x55)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BORDER = RGBColor(0x99, 0xA8, 0xBE)

FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"

REPO_DIR = Path(__file__).parent
IMAGES_DIR = REPO_DIR / "images"
OUTPUT = REPO_DIR / "presentation.pptx"

PLACEHOLDER_TAG = "[placeholder]"


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_solid(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _no_line(shape) -> None:
    shape.line.fill.background()


def _set_paragraph(
    p,
    text: str,
    *,
    size: int = 16,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor | None = None,
    font: str = FONT_BODY,
    align: PP_ALIGN | None = None,
    level: int = 0,
) -> None:
    p.text = text
    if align is not None:
        p.alignment = align
    p.level = level
    for run in p.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        if color is not None:
            run.font.color.rgb = color


def _format_runs(
    p,
    *,
    size: int = 16,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor | None = None,
    font: str = FONT_BODY,
) -> None:
    for run in p.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        if color is not None:
            run.font.color.rgb = color


def _add_title_bar(slide, title: str) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.85))
    _set_solid(bar, COLOR_PRIMARY)
    _no_line(bar)
    tf = bar.text_frame
    tf.margin_left = Inches(0.45)
    tf.margin_right = Inches(0.45)
    tf.margin_top = Inches(0.12)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_paragraph(
        tf.paragraphs[0],
        title,
        size=26,
        bold=True,
        color=COLOR_TITLE_BAR_TEXT,
        font=FONT_TITLE,
    )


def _add_footer_band(slide, text: str = "Cargo-to-Door | Draft v0.1") -> None:
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        SLIDE_H - Inches(0.35),
        SLIDE_W,
        Inches(0.35),
    )
    _set_solid(band, COLOR_PRIMARY_LIGHT)
    _no_line(band)
    tf = band.text_frame
    tf.margin_left = Inches(0.45)
    tf.margin_right = Inches(0.45)
    tf.margin_top = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_paragraph(
        tf.paragraphs[0],
        text,
        size=10,
        color=COLOR_SUBTLE,
        font=FONT_BODY,
    )


def _add_text_box(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    return box, tf


def _add_picture_with_border(slide, image_path: Path, x, y, w, h):
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x - Emu(9525), y - Emu(9525),
                                    w + Emu(19050), h + Emu(19050))
    border.fill.background()
    border.line.color.rgb = COLOR_BORDER
    border.line.width = Pt(0.75)
    pic = slide.shapes.add_picture(str(image_path), x, y, width=w, height=h)
    return pic


def _add_picture_fit(slide, image_path: Path, x, y, max_w, max_h):
    """Add a picture preserving aspect ratio, fitting inside (max_w, max_h)
    centered in that box. Returns the picture shape."""
    from PIL import Image  # ships with python-pptx via Pillow

    with Image.open(str(image_path)) as im:
        iw, ih = im.size
    aspect_img = iw / ih
    aspect_box = max_w / max_h
    if aspect_img >= aspect_box:
        w = max_w
        h = int(max_w / aspect_img)
    else:
        h = max_h
        w = int(max_h * aspect_img)
    x_centered = x + (max_w - w) // 2
    y_centered = y + (max_h - h) // 2

    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_centered - Emu(9525),
                                    y_centered - Emu(9525), w + Emu(19050), h + Emu(19050))
    border.fill.background()
    border.line.color.rgb = COLOR_BORDER
    border.line.width = Pt(0.75)
    pic = slide.shapes.add_picture(str(image_path), x_centered, y_centered,
                                   width=w, height=h)
    return pic


def add_title_slide(prs: Presentation, title: str, subtitle: str, byline: str) -> None:
    slide = _blank(prs)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(2.6))
    _set_solid(bar, COLOR_PRIMARY)
    _no_line(bar)

    _, tf = _add_text_box(slide, Inches(0.8), Inches(0.7), SLIDE_W - Inches(1.6), Inches(1.6))
    _set_paragraph(tf.paragraphs[0], title, size=54, bold=True,
                   color=COLOR_TITLE_BAR_TEXT, font=FONT_TITLE)

    _, tf = _add_text_box(slide, Inches(0.8), Inches(2.0), SLIDE_W - Inches(1.6), Inches(0.7))
    _set_paragraph(tf.paragraphs[0], subtitle, size=22, italic=True,
                   color=COLOR_TITLE_BAR_TEXT, font=FONT_TITLE)

    _, tf = _add_text_box(slide, Inches(0.8), Inches(3.4), SLIDE_W - Inches(1.6), Inches(3.0))
    _set_paragraph(tf.paragraphs[0],
                   "A building-integrated logistics platform for new high-rises and "
                   "historic retrofits.",
                   size=20, color=COLOR_BODY_TEXT, font=FONT_BODY)

    p = tf.add_paragraph()
    _set_paragraph(p, "", size=10)
    p2 = tf.add_paragraph()
    _set_paragraph(p2, byline, size=16, italic=True, color=COLOR_PLACEHOLDER, font=FONT_BODY)


def add_section_divider(prs: Presentation, label: str) -> None:
    slide = _blank(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_solid(bg, COLOR_DIVIDER_BG)
    _no_line(bg)

    _, tf = _add_text_box(slide, Inches(0.5), Inches(3.0),
                          SLIDE_W - Inches(1.0), Inches(1.5))
    _set_paragraph(tf.paragraphs[0], label, size=54, bold=True,
                   color=COLOR_WHITE, font=FONT_TITLE, align=PP_ALIGN.CENTER)


def add_bullets_slide(
    prs: Presentation,
    title: str,
    sections: Sequence[tuple[str, Sequence[str]]],
    *,
    callout: str | None = None,
) -> None:
    """sections = [(heading, [bullets]), ...]. Headings can be empty string."""
    slide = _blank(prs)
    _add_title_bar(slide, title)

    content_y = Inches(1.1)
    content_h = SLIDE_H - Inches(1.6)
    _, tf = _add_text_box(slide, Inches(0.5), content_y,
                          SLIDE_W - Inches(1.0), content_h)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    first = True
    for heading, bullets in sections:
        if heading:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            _set_paragraph(p, heading, size=18, bold=True, color=COLOR_PRIMARY,
                           font=FONT_BODY)
        for b in bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            is_placeholder = PLACEHOLDER_TAG in b
            display = b.replace(PLACEHOLDER_TAG + " ", "").replace(PLACEHOLDER_TAG, "")
            _set_paragraph(
                p,
                "• " + display,
                size=15,
                color=COLOR_PLACEHOLDER if is_placeholder else COLOR_BODY_TEXT,
                italic=is_placeholder,
                level=0,
                font=FONT_BODY,
            )

    if callout:
        callout_h = Inches(0.9)
        callout_y = SLIDE_H - Inches(0.35) - callout_h - Inches(0.1)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), callout_y,
            SLIDE_W - Inches(1.0), callout_h,
        )
        _set_solid(box, COLOR_PRIMARY_LIGHT)
        box.line.color.rgb = COLOR_PRIMARY
        box.line.width = Pt(0.75)
        tf2 = box.text_frame
        tf2.margin_left = Inches(0.25)
        tf2.margin_right = Inches(0.25)
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf2.word_wrap = True
        _set_paragraph(tf2.paragraphs[0], callout, size=14, italic=True,
                       color=COLOR_PRIMARY, font=FONT_BODY)

    _add_footer_band(slide)


def add_content_image_slide(
    prs: Presentation,
    title: str,
    bullets: Sequence[str],
    image_path: Path,
) -> None:
    slide = _blank(prs)
    _add_title_bar(slide, title)

    bullets_x = Inches(0.5)
    bullets_y = Inches(1.1)
    bullets_w = Inches(6.5)
    bullets_h = SLIDE_H - Inches(1.5)

    _, tf = _add_text_box(slide, bullets_x, bullets_y, bullets_w, bullets_h)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        is_placeholder = PLACEHOLDER_TAG in b
        display = b.replace(PLACEHOLDER_TAG + " ", "").replace(PLACEHOLDER_TAG, "")
        _set_paragraph(
            p,
            "• " + display,
            size=15,
            color=COLOR_PLACEHOLDER if is_placeholder else COLOR_BODY_TEXT,
            italic=is_placeholder,
            font=FONT_BODY,
        )

    img_x = Inches(7.2)
    img_y = Inches(1.1)
    img_w = Inches(5.7)
    img_h = SLIDE_H - Inches(1.5)
    _add_picture_fit(slide, image_path, img_x, img_y, img_w, img_h)
    _add_footer_band(slide)


def add_two_image_slide(
    prs: Presentation,
    title: str,
    bullets: Sequence[str],
    img_left: Path,
    img_right: Path,
    third_img: Path | None = None,
    third_caption: str = "",
) -> None:
    slide = _blank(prs)
    _add_title_bar(slide, title)

    bullets_y = Inches(1.0)
    bullets_h = Inches(1.6)
    _, tf = _add_text_box(slide, Inches(0.5), bullets_y,
                          SLIDE_W - Inches(1.0), bullets_h)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        is_placeholder = PLACEHOLDER_TAG in b
        display = b.replace(PLACEHOLDER_TAG + " ", "").replace(PLACEHOLDER_TAG, "")
        _set_paragraph(
            p,
            "• " + display,
            size=13,
            color=COLOR_PLACEHOLDER if is_placeholder else COLOR_BODY_TEXT,
            italic=is_placeholder,
        )

    img_top_y = Inches(2.75)
    img_h = SLIDE_H - Inches(0.4) - img_top_y

    if third_img is None:
        col_w = (SLIDE_W - Inches(1.5)) // 2
        _add_picture_fit(slide, img_left, Inches(0.5), img_top_y, col_w, img_h)
        _add_picture_fit(slide, img_right,
                         Inches(0.5) + col_w + Inches(0.5),
                         img_top_y, col_w, img_h)
    else:
        col_w = (SLIDE_W - Inches(1.5)) * 2 // 5
        third_w = SLIDE_W - Inches(1.5) - 2 * col_w - Inches(0.5)
        _add_picture_fit(slide, img_left, Inches(0.5), img_top_y, col_w, img_h)
        _add_picture_fit(slide, img_right,
                         Inches(0.5) + col_w + Inches(0.25),
                         img_top_y, col_w, img_h)
        third_x = Inches(0.5) + 2 * col_w + Inches(0.75)
        third_h = img_h - Inches(0.6)
        _add_picture_fit(slide, third_img, third_x, img_top_y, third_w, third_h)
        cap_y = img_top_y + third_h + Inches(0.05)
        _, cap_tf = _add_text_box(slide, third_x, cap_y, third_w, Inches(0.5))
        cap_tf.vertical_anchor = MSO_ANCHOR.TOP
        _set_paragraph(cap_tf.paragraphs[0], third_caption, size=10, italic=True,
                       color=COLOR_SUBTLE, align=PP_ALIGN.CENTER)

    _add_footer_band(slide)


def add_table_slide(
    prs: Presentation,
    title: str,
    headers: Sequence[str],
    rows: Sequence[Sequence[str]],
    *,
    placeholder_cells: Iterable[tuple[int, int]] | None = None,
    col_weights: Sequence[float] | None = None,
    preface_bullets: Sequence[str] | None = None,
    notes_bullets: Sequence[str] | None = None,
    body_font_size: int = 11,
) -> None:
    slide = _blank(prs)
    _add_title_bar(slide, title)
    placeholder_set = set(placeholder_cells or ())

    table_x = Inches(0.5)
    table_y = Inches(1.1)
    table_w = SLIDE_W - Inches(1.0)
    available_h = SLIDE_H - Inches(1.5) - table_y

    preface_h = Inches(0)
    if preface_bullets:
        preface_h = Inches(0.4 + 0.32 * len(preface_bullets))
        _, tf = _add_text_box(slide, table_x, table_y, table_w, preface_h)
        for i, b in enumerate(preface_bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            is_placeholder = PLACEHOLDER_TAG in b
            display = b.replace(PLACEHOLDER_TAG + " ", "").replace(PLACEHOLDER_TAG, "")
            _set_paragraph(
                p, "• " + display, size=13,
                color=COLOR_PLACEHOLDER if is_placeholder else COLOR_BODY_TEXT,
                italic=is_placeholder,
            )
        table_y = table_y + preface_h + Inches(0.1)
        available_h = SLIDE_H - Inches(1.5) - table_y

    notes_h = Inches(0)
    if notes_bullets:
        notes_h = Inches(0.4 + 0.32 * len(notes_bullets))
        available_h = available_h - notes_h - Inches(0.1)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_h = available_h
    table_shape = slide.shapes.add_table(n_rows, n_cols, table_x, table_y, table_w, table_h)
    table = table_shape.table

    if col_weights:
        total = sum(col_weights)
        widths = [int(table_w * w / total) for w in col_weights]
        delta = int(table_w) - sum(widths)
        widths[-1] += delta
        for i, w in enumerate(widths):
            table.columns[i].width = w

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        cell.text = h
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            _format_runs(p, size=12, bold=True, color=COLOR_WHITE, font=FONT_BODY)

    for r, row in enumerate(rows, start=1):
        zebra = (r % 2 == 0)
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            if zebra:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_PRIMARY_LIGHT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_WHITE
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.TOP
            is_placeholder_cell = (r - 1, c) in placeholder_set or PLACEHOLDER_TAG in value
            display = value.replace(PLACEHOLDER_TAG + " ", "").replace(PLACEHOLDER_TAG, "")
            cell.text = display
            for p in cell.text_frame.paragraphs:
                _format_runs(
                    p, size=body_font_size,
                    color=COLOR_PLACEHOLDER if is_placeholder_cell else COLOR_BODY_TEXT,
                    italic=is_placeholder_cell,
                )

    if notes_bullets:
        notes_y = table_y + table_h + Inches(0.1)
        _, tf = _add_text_box(slide, table_x, notes_y, table_w, notes_h)
        for i, b in enumerate(notes_bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            is_placeholder = PLACEHOLDER_TAG in b
            display = b.replace(PLACEHOLDER_TAG + " ", "").replace(PLACEHOLDER_TAG, "")
            _set_paragraph(
                p, "• " + display, size=12,
                color=COLOR_PLACEHOLDER if is_placeholder else COLOR_BODY_TEXT,
                italic=is_placeholder,
            )

    _add_footer_band(slide)


def _add_box(slide, x, y, w, h, label, *, fill=None, line_color=None,
             font_size=11, bold=True, font_color=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    if fill is not None:
        _set_solid(box, fill)
    if line_color is not None:
        box.line.color.rgb = line_color
        box.line.width = Pt(0.75)
    else:
        box.line.color.rgb = COLOR_PRIMARY
        box.line.width = Pt(0.75)
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    _set_paragraph(
        tf.paragraphs[0], label, size=font_size, bold=bold,
        color=font_color if font_color is not None else COLOR_PRIMARY,
        align=PP_ALIGN.CENTER, font=FONT_BODY,
    )
    return box


def _add_arrow(slide, x1, y1, x2, y2, *, dashed: bool = False, color=None):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color if color is not None else COLOR_PRIMARY
    conn.line.width = Pt(1.25)
    ln = conn.line._get_or_add_ln()
    tail_end = etree.SubElement(ln, qn("a:tailEnd"))
    tail_end.set("type", "triangle")
    tail_end.set("w", "med")
    tail_end.set("len", "med")
    if dashed:
        dash = etree.SubElement(ln, qn("a:prstDash"))
        dash.set("val", "dash")
    return conn


def add_product_overview_flow_slide(prs: Presentation, title: str) -> None:
    """Slide S10: LR end-to-end flow with sky and street branches feeding the core."""
    slide = _blank(prs)
    _add_title_bar(slide, title)

    box_w = Inches(1.55)
    box_h = Inches(0.7)
    gap_x = Inches(0.13)

    sky_y = Inches(1.5)
    sky_nodes = [
        "Drone", "Rooftop\nFlight Deck", "Catchment\nFunnel + S-Trap",
        "Sorting\nShaft", "Smart\nDumbwaiter", "Smart\nLocker", "Resident",
    ]
    sky_boxes = []
    x = Inches(0.5)
    for label in sky_nodes:
        b = _add_box(slide, x, sky_y, box_w, box_h, label,
                     fill=COLOR_PRIMARY_LIGHT, font_size=10)
        sky_boxes.append(b)
        x += box_w + gap_x

    for i in range(len(sky_boxes) - 1):
        a = sky_boxes[i]
        b = sky_boxes[i + 1]
        _add_arrow(slide,
                   a.left + a.width, a.top + a.height // 2,
                   b.left, b.top + b.height // 2)

    street_y = Inches(3.4)
    street_labels = ["Micro-Rover", "Kinetic\nDocking Hatch", "Robotic-Arm\nVestibule"]
    street_boxes = []
    x = Inches(0.5)
    for label in street_labels:
        b = _add_box(slide, x, street_y, box_w, box_h, label,
                     fill=COLOR_PRIMARY_LIGHT, font_size=10)
        street_boxes.append(b)
        x += box_w + gap_x

    for i in range(len(street_boxes) - 1):
        a = street_boxes[i]
        b = street_boxes[i + 1]
        _add_arrow(slide,
                   a.left + a.width, a.top + a.height // 2,
                   b.left, b.top + b.height // 2)

    arm = street_boxes[-1]
    sort_box = sky_boxes[3]
    _add_arrow(slide,
               arm.left + arm.width // 2, arm.top,
               sort_box.left + sort_box.width // 2, sort_box.top + sort_box.height)

    cloud_y = Inches(5.3)
    cloud_box = _add_box(
        slide, Inches(4.5), cloud_y, Inches(4.3), Inches(0.7),
        "Cloud Orchestration\n(Fleet Manager • Routing • Resident App)",
        fill=COLOR_PRIMARY, font_color=COLOR_WHITE, font_size=11,
    )

    for target in (sky_boxes[0], street_boxes[0], sky_boxes[5]):
        _add_arrow(
            slide,
            cloud_box.left + cloud_box.width // 2,
            cloud_box.top,
            target.left + target.width // 2,
            target.top + target.height,
            dashed=True,
        )

    _, tf = _add_text_box(slide, Inches(0.5), Inches(6.3),
                          SLIDE_W - Inches(1.0), Inches(0.7))
    _set_paragraph(tf.paragraphs[0],
                   "Two intake modalities (sky + street) converge on a single internal "
                   "sorting core distributing to resident smart-lockers — orchestrated "
                   "by the cloud layer.",
                   size=12, italic=True, color=COLOR_SUBTLE)

    _add_footer_band(slide)


def add_architecture_slide(prs: Presentation, title: str) -> None:
    """Slide S12: layered subsystem swimlanes (Airspace, Rooftop, Vertical Core,
    Floor, Street, Cloud)."""
    slide = _blank(prs)
    _add_title_bar(slide, title)

    lanes = [
        ("Airspace",       ["Drone Fleet", "Airspace Control / Geofence"]),
        ("Rooftop Hub",    ["Flight Deck", "Catchment Funnel + S-Trap"]),
        ("Vertical Core",  ["Sorting Shaft", "Smart Dumbwaiter"]),
        ("Floor Layer",    ["Automatic Floor Switcher", "Smart-Locker Array"]),
        ("Street Layer",   ["Micro-Rover", "Kinetic Docking Hatch", "Vestibule + Robotic Arm"]),
        ("Cloud",          ["Fleet Manager", "Routing & Scheduling", "Resident App"]),
    ]

    top = Inches(1.05)
    bottom = SLIDE_H - Inches(0.5)
    total_h = bottom - top
    n = len(lanes)
    gap = Inches(0.08)
    lane_h = (total_h - gap * (n - 1)) // n

    lane_x = Inches(0.5)
    lane_w = SLIDE_W - Inches(1.0)
    label_w = Inches(1.8)

    inner_left = lane_x + label_w + Inches(0.15)
    inner_right = lane_x + lane_w - Inches(0.15)
    inner_w_total = inner_right - inner_left

    for i, (label, nodes) in enumerate(lanes):
        y = top + i * (lane_h + gap)
        lane_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         lane_x, y, lane_w, lane_h)
        _set_solid(lane_bg, COLOR_PRIMARY_LIGHT)
        lane_bg.line.color.rgb = COLOR_PRIMARY
        lane_bg.line.width = Pt(0.5)

        label_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           lane_x, y, label_w, lane_h)
        _set_solid(label_box, COLOR_PRIMARY)
        _no_line(label_box)
        tf = label_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set_paragraph(tf.paragraphs[0], label, size=12, bold=True,
                       color=COLOR_WHITE, align=PP_ALIGN.CENTER, font=FONT_BODY)

        m = len(nodes)
        node_gap = Inches(0.15)
        node_w = (inner_w_total - node_gap * (m - 1)) // m
        node_h = lane_h - Inches(0.18)
        node_y = y + Inches(0.09)
        nx = inner_left
        for node in nodes:
            _add_box(slide, nx, node_y, node_w, node_h, node,
                     fill=COLOR_WHITE, font_size=10)
            nx += node_w + node_gap

    _add_footer_band(slide)


def add_gantt_slide(
    prs: Presentation,
    title: str,
    wp_rows: Sequence[tuple[str, int, int, str]],
    *,
    months: int = 24,
) -> None:
    """wp_rows = [(name, start_month_1indexed, length_months, owner), ...]"""
    slide = _blank(prs)
    _add_title_bar(slide, title)

    n_rows = len(wp_rows) + 1
    n_cols = months + 2
    table_x = Inches(0.4)
    table_y = Inches(1.1)
    table_w = SLIDE_W - Inches(0.8)
    table_h = SLIDE_H - Inches(1.5) - table_y

    table_shape = slide.shapes.add_table(n_rows, n_cols, table_x, table_y, table_w, table_h)
    table = table_shape.table

    wp_w = Inches(1.7)
    owner_w = Inches(1.1)
    month_total = table_w - wp_w - owner_w
    month_w = month_total // months

    widths = [wp_w] + [month_w] * months + [owner_w]
    delta = int(table_w) - sum(widths)
    widths[1 + months // 2] += delta
    for i, w in enumerate(widths):
        table.columns[i].width = w

    headers = ["Work Package"] + [str(m + 1) for m in range(months)] + ["Owner"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        cell.text = h
        cell.margin_left = Inches(0.03)
        cell.margin_right = Inches(0.03)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            _format_runs(p, size=9, bold=True, color=COLOR_WHITE)
            p.alignment = PP_ALIGN.CENTER

    for r, (name, start, length, owner) in enumerate(wp_rows, start=1):
        zebra = (r % 2 == 0)

        cell = table.cell(r, 0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY_LIGHT if zebra else COLOR_WHITE
        cell.text = name
        cell.margin_left = Inches(0.05)
        cell.margin_right = Inches(0.05)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            _format_runs(p, size=10, bold=True, color=COLOR_PRIMARY)

        for m in range(months):
            mc = table.cell(r, 1 + m)
            in_bar = (m + 1) >= start and (m + 1) < (start + length)
            mc.fill.solid()
            if in_bar:
                mc.fill.fore_color.rgb = COLOR_PRIMARY
            else:
                mc.fill.fore_color.rgb = COLOR_PRIMARY_LIGHT if zebra else COLOR_WHITE
            mc.text = ""

        oc = table.cell(r, n_cols - 1)
        oc.fill.solid()
        oc.fill.fore_color.rgb = COLOR_PRIMARY_LIGHT if zebra else COLOR_WHITE
        is_placeholder = PLACEHOLDER_TAG in owner
        display = owner.replace(PLACEHOLDER_TAG + " ", "").replace(PLACEHOLDER_TAG, "")
        oc.text = display if display else "[name]"
        oc.margin_left = Inches(0.05)
        oc.margin_right = Inches(0.05)
        oc.margin_top = Inches(0.04)
        oc.margin_bottom = Inches(0.04)
        oc.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in oc.text_frame.paragraphs:
            _format_runs(
                p, size=9,
                color=COLOR_PLACEHOLDER if is_placeholder else COLOR_BODY_TEXT,
                italic=is_placeholder,
            )
            p.alignment = PP_ALIGN.CENTER

    _add_footer_band(slide)


def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(
        prs,
        "Cargo-to-Door",
        "Autonomous Last-Meter Delivery for Vertical Cities",
        "[placeholder] Presented by: Team Name — PI, Co-PI, Members — Institution — Date",
    )

    add_bullets_slide(
        prs,
        "Agenda",
        [
            ("1. Background", [
                "Motivation & Problem Definition",
                "Literature Survey: state of the art, our proposed solution, breakthrough technologies",
            ]),
            ("2. The Team & Innovative Product Development", [
                "Team & capabilities, requirements, architecture, design, implementation, testing",
            ]),
            ("3. Roadmap / Project Plan", [
                "Milestones, deliverables, work packages, timeline",
            ]),
            ("4. Commercialization & Business Model", [
                "Market, competition, budget, risks",
            ]),
        ],
    )

    add_section_divider(prs, "Part 1 — Background")

    add_bullets_slide(
        prs,
        "Motivation & Problem Definition",
        [
            ("", [
                "Last-mile is the broken link — ~30–50% of total e-commerce logistics cost, plus urban congestion and emissions.",
                "The \"last-meter\" problem: even at the lobby, the trip to the apartment door is manual, slow, and unreliable.",
                "Failed deliveries & porch piracy drive returns, theft, and customer churn.",
                "Vertical density breaks current solutions: drones cannot land on balconies, rovers cannot ride elevators, couriers cannot scale to 50-storey towers.",
                "Historic & protected buildings cannot be torn open for new infrastructure — yet residents demand the same service.",
            ]),
        ],
        callout="Goal: a continuous, autonomous, weatherproof path from sky (or sidewalk) "
                "to the resident's door — for new construction AND historic retrofits.",
    )

    add_table_slide(
        prs,
        "Literature Survey — Past Solutions / State of the Art",
        ["Approach", "Examples", "Strengths", "Limitations"],
        [
            ["Drone-to-doorstep", "Wing, Zipline, Matternet",
             "Bypasses ground congestion",
             "Needs landing pad; weather-sensitive; one parcel per flight; not vertical-city friendly"],
            ["Sidewalk micro-rovers", "Starship, Serve, Coco",
             "Low-cost, social-friendly",
             "Stuck at the lobby; no vertical access; theft risk"],
            ["Parcel-locker arrays", "Amazon Hub, InPost",
             "Secure, asynchronous pickup",
             "Still requires resident trip; ground-floor only"],
            ["Pneumatic / vacuum tubes", "Hospital legacy; Hyperloop-cargo concepts",
             "Weather-immune, high throughput",
             "Single-routing; expensive retrofit; small payload"],
            ["In-building robotics", "Relay, Savioke",
             "Door-to-door inside building",
             "Needs elevator integration; no exterior link"],
        ],
        col_weights=[1.5, 2.0, 2.0, 3.0],
        notes_bullets=[
            "Gap: no end-to-end system spans airspace -> building skin -> vertical core -> resident door, and none retrofits gracefully onto protected architecture.",
        ],
        body_font_size=10,
    )

    add_content_image_slide(
        prs,
        "Proposed Solution & Innovation",
        [
            "Two parallel tracks, one platform.",
            "New builds — Rooftop Flight Deck + internal vertical shaft. Drones hover and winch-drop into an aerodynamic catchment funnel — no landing required.",
            "Historic builds — Exo-Logistics Spine: modular exterior pneumatic-decelerator tube clipped to the facade, feeding automatic floor-switchers inside.",
            "Passive-logistics philosophy: gravity chutes, motorized rollers, mechanical diverters and claws replace expensive robotic arms wherever possible.",
            "Multi-modal intake: same vertical core accepts both aerial drone and sidewalk micro-rover deliveries.",
        ],
        IMAGES_DIR / "new-building-rooftop.png",
    )

    add_bullets_slide(
        prs,
        "Breakthrough Technologies Required",
        [
            ("", [
                "Hover-and-winch package drop — drone holds altitude Z while lowering a parcel; controller compensates for wind and tether sway.",
                "Aerodynamic catchment funnel + S-Trap drainage — captures off-center drops, lets rain through, blocks rain from following the parcel.",
                "Pneumatic decelerator tube — variable-pressure exterior shaft that slows a free-falling parcel from terminal velocity to safe handoff speed.",
                "Automatic floor-switcher — passive-actuated diverter that routes packages into the correct resident's intake chamber.",
                "Micro-rover <-> kinetic docking hatch — standardized mechanical & data interface so any six-wheeled rover can offload into the building.",
                "Building-level airspace orchestration — scheduling, geofencing, conflict-resolution for multiple drones over a single rooftop.",
            ]),
        ],
    )

    add_section_divider(prs, "Part 2 — Team & Innovative Product")

    add_table_slide(
        prs,
        "The Team & Capabilities",
        ["Role", "Person", "Capability brought"],
        [
            ["Principal Investigator", "[placeholder] [name]",  "Systems engineering, autonomous vehicles"],
            ["Mechanical Lead",        "[placeholder] [name]",  "Aerodynamic structures, pneumatic systems"],
            ["Controls & Robotics",    "[placeholder] [name]",  "Drone control, robotic manipulation"],
            ["Software & Cloud",       "[placeholder] [name]",  "Fleet orchestration, edge computing"],
            ["Architecture & Civil",   "[placeholder] [name]",  "High-rise integration, historic-preservation retrofits"],
            ["Business & Regulatory",  "[placeholder] [name]",  "Go-to-market, FAA / municipal liaison"],
        ],
        col_weights=[1.6, 1.6, 4.0],
        preface_bullets=["[placeholder] Fill with real names, affiliations, and CVs before submission."],
        notes_bullets=[
            "Institutional capabilities: wind-tunnel access, drop-test rig, full-scale mockup hall, partnership LOIs with [placeholder] [developer] and [placeholder] [historic-district BID].",
        ],
        body_font_size=12,
    )

    add_product_overview_flow_slide(prs, "Innovative Product / Process — System Overview")

    add_bullets_slide(
        prs,
        "Requirements Definition & Analysis",
        [
            ("Functional", [
                "Throughput: [placeholder] >= 60 pkg/hr per building at peak",
                "Payload envelope: [placeholder] <= 5 kg, <= 40 x 30 x 30 cm",
                "End-to-end latency drone-arrival -> locker: [placeholder] < 4 min",
            ]),
            ("Non-functional", [
                "Weatherproof to [placeholder] IP65 equivalent; operable in [placeholder] wind <= 12 m/s, rain <= 25 mm/h",
                "Acoustic limit at facade: [placeholder] <= 55 dBA at 10 m",
                "Fire-rating of internal shafts: [placeholder] >= 2-hour",
            ]),
            ("Regulatory / Safety", [
                "FAA Part 107 (or local equivalent) drone operations",
                "Local building, fire, and (for retrofits) landmark-commission codes",
            ]),
            ("UX", [
                "Resident retrieval at locker: [placeholder] < 30 s; mobile app notify within [placeholder] 10 s of locker close",
            ]),
        ],
    )

    add_architecture_slide(prs, "High-Level System Architecture")

    add_content_image_slide(
        prs,
        "Design — Rooftop Flight Deck",
        [
            "Stable delivery zone with perforated wind-break louvers to tame rooftop turbulence.",
            "Aerodynamic catchment funnel sized for off-center hover drops.",
            "Secure storage lockers for high-value or oversized parcels awaiting retrieval.",
            "Hover-winch protocol: drone holds altitude Z, lowers parcel through funnel mouth, releases tether — never lands.",
            "Flight-path corridor geofenced into building airspace; multiple drones serialized on approach.",
        ],
        IMAGES_DIR / "new-building-rooftop.png",
    )

    add_content_image_slide(
        prs,
        "Design — Drop Funnel & S-Trap Drainage",
        [
            "Drone drop aperture with Teflon-grated inner skin — parcel slides, water passes through.",
            "U-bend / S-Trap mechanism — passive drainage; rainwater diverts to the storm drain, parcel continues down the Momentum Carry Zone.",
            "Dry interior zone opens onto the autonomous internal sorting shaft and roller conveyor.",
            "Solves the rain-in-shaft constraint from content.md — no powered seals, no moving water doors.",
        ],
        IMAGES_DIR / "new-building-get-cargo.png",
    )

    add_content_image_slide(
        prs,
        "Design — Internal Sorting Core",
        [
            "Stainless-steel gravity chute delivers parcels from the rooftop core to floor-level receiving bays.",
            "Motorized roller conveyor + diverter flaps route each parcel to its destination column.",
            "Smart dumbwaiter (open-front car in a dedicated shaft) lifts/lowers to the resident's floor.",
            "Embodies the passive-logistics principle from content.md: gravity + rollers + mechanical diverters in place of robotic arms.",
        ],
        IMAGES_DIR / "new-building-seperator.png",
    )

    add_content_image_slide(
        prs,
        "Design — Resident Smart-Locker Array",
        [
            "Integrated smart-locker array built into the apartment-floor hallway wall — wood facade, glass + tech inserts.",
            "Sliding secure logistics core behind the lockers connects directly to the vertical shaft.",
            "Resident gets a push notification (\"Package received: Unit 4C, Compartment L16\"), authenticates at the screen, locker opens.",
            "Inset on the source drawing details the automated dispatch path: internal chute -> conveyor sorting -> dumbwaiter shaft.",
        ],
        IMAGES_DIR / "new-building-smart-locker.png",
    )

    add_two_image_slide(
        prs,
        "Design — Retrofit \"Exo-Logistics Spine\"",
        [
            "Exterior modular tube clamped to the facade of historic buildings — no structural intervention to the protected envelope.",
            "Pneumatic decelerator inside the tube progressively brakes a free-falling parcel.",
            "Automatic switcher mechanism (right) diverts the parcel into each floor's resident intake chamber.",
            "Rooftop drone docking station caps the spine; a drone lowers the parcel into the tube mouth.",
        ],
        IMAGES_DIR / "retrofit-pipe-desing-outside.png",
        IMAGES_DIR / "retrofit-pipe-design-switcher.png",
    )

    add_two_image_slide(
        prs,
        "Design — Ground-Level Micro-Rover Interface",
        [
            "Segregated micro-logistics lane along the sidewalk; rover approaches the kinetic docking hatch in the facade (left).",
            "Inside the vestibule (middle), a ceiling-mounted precision robotic arm lifts parcels off the rover's open cargo bin onto a motorized roller conveyor.",
            "Retrofit equivalent (right): street-level micro-vestibule with a telescoping claw + motorized lift-tray.",
        ],
        IMAGES_DIR / "new-building-rover.png",
        IMAGES_DIR / "new-building-rover-get.png",
        third_img=IMAGES_DIR / "retrofit-claw.png",
        third_caption="Retrofit equivalent — telescoping claw vestibule",
    )

    add_table_slide(
        prs,
        "Implementation & Testing",
        ["KPI", "Target"],
        [
            ["Drop capture success rate",        "[placeholder] >= 99.5%"],
            ["Mean time, rooftop -> locker",     "[placeholder] <= 3 min"],
            ["Facade acoustic level",            "[placeholder] <= 55 dBA @ 10 m"],
            ["Rain ingress past S-Trap",         "[placeholder] 0 ml/h @ 25 mm/h rainfall"],
            ["Rover handoff cycle time",         "[placeholder] <= 45 s"],
        ],
        col_weights=[3.0, 2.5],
        preface_bullets=[
            "Subscale rooftop funnel + instrumented drop rig — characterize capture cone vs. wind speed.",
            "Single-floor mockup of internal sorting + smart-locker — measure throughput and jam rate.",
            "Exo-Spine pilot section (3 storeys) on a non-occupied historic structure — validate pneumatic deceleration and switcher routing.",
            "Rover <-> hatch integration with a partner micro-rover OEM.",
        ],
        body_font_size=12,
    )

    add_section_divider(prs, "Part 3 — Roadmap / Project Plan")

    add_table_slide(
        prs,
        "Roadmap / Project Plan",
        ["Phase", "Months", "Focus", "Exit gate"],
        [
            ["P1 — R&D",                          "1–4",   "Concept refinement, simulation, safety case",        "Design review #1"],
            ["P2 — Subsystem prototypes",         "5–10",  "Funnel, S-Trap, switcher, hatch, rover dock",        "Bench-test sign-off"],
            ["P3 — Integrated single-floor pilot","11–16", "New-build mockup + Exo-Spine 3-storey rig",          "End-to-end KPI demo"],
            ["P4 — Certification & approvals",    "17–20", "FAA waiver, fire & building code, landmark commission","Permits granted"],
            ["P5 — Commercial pilot",             "21–24", "One new-build tower + one retrofit block",           "Paying first customer"],
        ],
        col_weights=[2.6, 1.0, 4.5, 2.4],
        preface_bullets=["[placeholder] 24-month phased plan — refine with real funding milestones."],
        body_font_size=11,
    )

    add_table_slide(
        prs,
        "Milestones, Deliverables & Work Packages",
        ["WP", "Title", "Key deliverables", "Lead", "Month"],
        [
            ["WP1", "Rooftop Hub",                "Flight Deck mock-up, hover-drop protocol spec", "[placeholder] [name]", "M8"],
            ["WP2", "Vertical Core",              "Sorting-shaft prototype, dumbwaiter integration","[placeholder] [name]","M12"],
            ["WP3", "Exo-Spine",                  "Pneumatic-tube section, automatic switcher",    "[placeholder] [name]", "M14"],
            ["WP4", "Rover Interface",            "Kinetic docking hatch + robotic-arm vestibule", "[placeholder] [name]", "M12"],
            ["WP5", "Software & Orchestration",   "Fleet manager, resident app, routing engine",   "[placeholder] [name]", "M16"],
            ["WP6", "Regulatory & Safety",        "FAA waiver, fire/building/landmark approvals",  "[placeholder] [name]", "M20"],
            ["WP7", "Pilot Deployment",           "One new-build + one retrofit operating",        "[placeholder] [name]", "M24"],
        ],
        col_weights=[0.6, 2.4, 4.5, 1.6, 0.7],
        preface_bullets=["[placeholder] Owners and exact months TBD."],
        body_font_size=10,
    )

    add_gantt_slide(
        prs,
        "Task Allocation & Timeline (24 months)",
        [
            ("WP1 Rooftop Hub",       1, 7,  "[placeholder] [name]"),
            ("WP2 Vertical Core",     2, 7,  "[placeholder] [name]"),
            ("WP3 Exo-Spine",         4, 7,  "[placeholder] [name]"),
            ("WP4 Rover Interface",   3, 5,  "[placeholder] [name]"),
            ("WP5 Software",          1, 12, "[placeholder] [name]"),
            ("WP6 Regulatory",       11, 6,  "[placeholder] [name]"),
            ("WP7 Pilot Deployment", 17, 4,  "[placeholder] [name]"),
        ],
    )

    add_section_divider(prs, "Part 4 — Commercialization & Business Model")

    add_bullets_slide(
        prs,
        "Market Analysis",
        [
            ("Target market", [
                "Luxury high-rise residential developers (new construction)",
                "Smart-city districts and master-planned communities",
                "Historic-district Business Improvement Districts commissioning the Exo-Spine retrofit",
                "Mixed-use towers with concierge service expectations",
            ]),
            ("Customer profile", [
                "Developer / building owner buys the infrastructure; residents are the end users",
                "Logistics providers (Amazon, FedEx, UPS, local couriers) are channel partners paying per-delivery",
            ]),
            ("Bureaucracy & approvals", [
                "FAA Part 107 + local airspace authority",
                "Building & fire code (UL listing for shafts)",
                "Landmark / historic-preservation commissions (retrofit only)",
            ]),
            ("Market size", [
                "[placeholder] TAM: $XX B global vertical-living logistics",
                "[placeholder] SAM: $X B Tier-1 cities with drone-friendly regs",
                "[placeholder] SOM (Year 5): $XXX M",
            ]),
            ("Marketing & growth strategy", [
                "Lighthouse pilot with one flagship developer -> case study -> master-developer agreements",
                "\"Logistics-ready building\" certification co-marketed with developers",
                "API partnerships with major carriers for instant scale on Day 1 of a building going live",
            ]),
        ],
    )

    add_table_slide(
        prs,
        "Competitor Analysis",
        ["Competitor", "Their play", "Where we win"],
        [
            ["Wing / Zipline / Matternet", "Drone-to-doorstep delivery",
             "We solve the last-meter (lobby -> door); they stop at the curb / yard"],
            ["Starship / Serve / Coco", "Sidewalk micro-rovers",
             "We integrate with their rovers via a standard hatch — we are the building-side complement, not a competitor"],
            ["Amazon Hub / InPost", "Ground-floor parcel lockers",
             "Our smart-lockers are on the resident's own floor, fed automatically — no lobby trip"],
            ["Relay / Savioke", "In-building delivery robots",
             "We bypass elevators with a dedicated shaft; lower opex, higher throughput"],
            ["Pneumatic-tube vendors", "Hospital-grade tubes",
             "We add aerial intake, weather handling, and resident-locker layer"],
        ],
        col_weights=[2.2, 2.5, 5.0],
        notes_bullets=[
            "Our competitive advantage: only end-to-end, building-integrated stack from airspace to apartment door; weatherproof passive shafts; retrofit option (Exo-Spine) unlocks the entire historic-building market; modality-agnostic (drones, rovers, human couriers).",
        ],
        body_font_size=10,
    )

    add_table_slide(
        prs,
        "Budget, Pricing & Monetization",
        ["Line", "Year 1", "Year 3"],
        [
            ["R&D / engineering",                  "[placeholder] $X M",  "[placeholder] $X M"],
            ["Hardware (per-building install)",    "[placeholder] $XXX k","[placeholder] $XXX k"],
            ["Software / cloud",                   "[placeholder] $X M",  "[placeholder] $X M"],
            ["Operations & support",               "[placeholder] $X M",  "[placeholder] $X M"],
        ],
        col_weights=[3.5, 1.5, 1.5],
        preface_bullets=["[placeholder] Cost structure — fill with real numbers."],
        notes_bullets=[
            "Pricing models: capex sale to developers / BIDs; SaaS orchestration per building per month; per-delivery fee to logistics partners; premium-resident subscription.",
            "[placeholder] Unit economics: payback per building ~XX months; gross margin at scale ~XX%.",
        ],
        body_font_size=12,
    )

    add_table_slide(
        prs,
        "Challenges & Business Risks",
        ["Risk", "Likelihood", "Impact", "Mitigation"],
        [
            ["Urban drone-airspace regulation slips", "High", "High",
             "Early FAA waiver process; partner with city pilot programs"],
            ["Public acceptance (noise, privacy, drones overhead)", "Med", "High",
             "Acoustic engineering at facade; opt-in resident comms; community design reviews"],
            ["Insurance & liability for sky-dropped parcels", "Med", "Med",
             "Capped payload, redundant tether, geofenced corridor, partner with logistics-insurance underwriter"],
            ["Landmark-commission rejection of Exo-Spine", "Med", "High",
             "Reversible / non-penetrating clamps; engage preservation architects from day 1"],
            ["Weather window too narrow (wind, ice, lightning)", "Med", "Med",
             "Hybrid sky + street modality — rovers cover the window where drones can't fly"],
            ["Hardware reliability (jammed switcher, frozen rollers)", "Med", "Med",
             "Passive-mechanical design philosophy; predictive maintenance via cloud telemetry"],
            ["Capital intensity / long sales cycle to developers", "High", "Med",
             "SaaS layer for recurring revenue; lighthouse-pilot case-study to compress sales cycle"],
        ],
        col_weights=[3.8, 1.0, 1.0, 4.5],
        body_font_size=10,
    )

    add_title_slide(
        prs,
        "Thank You",
        "Questions & Discussion",
        "[placeholder] contact: team@cargo-to-door.example | repo: github.com/<org>/cargo-to-door-design",
    )

    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
